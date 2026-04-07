"""MrMemory Acqui-Hire Pitch Deck Builder"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x0A, 0x0A, 0x0B)
BG2 = RGBColor(0x11, 0x11, 0x13)
BG3 = RGBColor(0x18, 0x18, 0x1B)
ACCENT = RGBColor(0x6D, 0x5C, 0xFF)
ACCENT2 = RGBColor(0x8B, 0x7A, 0xFF)
WHITE = RGBColor(0xFA, 0xFA, 0xFA)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
DARK_GRAY = RGBColor(0x71, 0x71, 0x7A)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BORDER = RGBColor(0x27, 0x27, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
    return txBox

def add_code_block(slide, left, top, width, height, code, font_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = ACCENT2
    p.font.name = "Cascadia Code"
    return shape

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, title_size=18, body_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_right = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = ACCENT2
    p.font.bold = True
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(body_size)
    p2.font.color.rgb = GRAY
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(8)
    return shape

# ============================================================
# SLIDE 1: Title / Hook
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide)
add_accent_line(slide, 0.8, 1.5, 2)
add_text(slide, 0.8, 1.7, 11, 1.5, "MrMemory", font_size=54, color=WHITE, bold=True)
add_text(slide, 0.8, 2.6, 11, 0.8, "Intelligent + Real-Time Memory Layer for Production Agents", font_size=28, color=ACCENT2)

add_text(slide, 0.8, 3.8, 11, 0.6, "One-line install. Agents remember forever — with auto-extraction, compression, self-edits, LangGraph-native persistence, and live multi-agent sync.", font_size=16, color=GRAY)

add_code_block(slide, 0.8, 4.7, 5.5, 1.1,
    'from mrmemory import AMR\namr = AMR("amr_sk_...")\namr.auto_remember(messages, sync=True)', font_size=14)

add_text(slide, 3, 6.2, 7, 0.4, "Darren R  •  Solo Founder  •  Shipped v0.4.1 March 2026", font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, 3, 6.6, 7, 0.4, "Acqui-Hire Discussion — March 2026", font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "The Agent Memory Problem", font_size=40, color=WHITE, bold=True)
add_text(slide, 0.8, 1.9, 10, 0.6, "In 2026, agents still forget across sessions or require heavy manual orchestration.", font_size=18, color=ACCENT2)

add_card(slide, 0.8, 2.8, 3.5, 2.0, "💸 Stuff the Prompt",
    "Shove everything in context.\nToken costs explode.\nContext windows overflow.")
add_card(slide, 4.7, 2.8, 3.5, 2.0, "🔧 Build Your Own Pipeline",
    "Roll a vector DB pipeline.\nWeeks of work.\nNot your core product.")
add_card(slide, 8.6, 2.8, 3.5, 2.0, "🤷 Skip Memory",
    "Skip memory entirely.\nAgents start from zero.\nUsers suffer.")

add_text(slide, 0.8, 5.2, 11, 1.0,
    "• LangGraph checkpointers = state snapshots only\n• Raw vector DBs = embeddings, no intelligence\n• No solution combines managed + intelligent + real-time",
    font_size=15, color=GRAY)

# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "MrMemory — The Complete Memory Layer", font_size=40, color=WHITE, bold=True)

add_card(slide, 0.8, 2.0, 3.7, 1.5, "🧠 Core API",
    "remember() — store with auto-embedding\nrecall() — semantic vector search\nshare() — multi-agent permissioned sharing")
add_card(slide, 4.8, 2.0, 3.7, 1.5, "✨ Intelligent Features",
    "auto_remember(messages) — LLM extraction,\ndedup, entity tagging, auto-TTL\ncompress() — group + summarize + metrics")
add_card(slide, 8.8, 2.0, 3.7, 1.5, "✏️ Self-Edit Tools",
    "update() — patch content, re-embed\ndelete_outdated() — bulk prune by age/tags\nmerge() — combine N memories into one")

add_card(slide, 0.8, 3.8, 5.6, 1.5, "⚡ Real-Time Multi-Agent",
    "WebSocket events for live memory sharing between agents.\nAgent A stores → Agent B gets it instantly.\nMemory.created, memory.updated, memory.deleted events.")
add_card(slide, 6.7, 3.8, 5.8, 1.5, "🔗 LangGraph Native",
    "MrMemoryCheckpointer — drop-in LangGraph persistence\nMrMemoryStore — cross-session memory for agents\npip install mrmemory[langchain] — one line integration")

add_text(slide, 0.8, 5.7, 11, 0.5, "All features ship today in mrmemory 0.4.1 (PyPI) and memorymr 0.4.1 (npm)", font_size=14, color=DARK_GRAY)

# ============================================================
# SLIDE 4: Technology & Moat
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "Core Technology & Moat", font_size=40, color=WHITE, bold=True)

add_code_block(slide, 0.8, 2.0, 7, 2.8,
    """┌─────────────┐     ┌──────────────┐     ┌─────────────┐
│  Your Agent  │────▶│  MrMemory    │────▶│   Qdrant    │
│  (SDK/REST)  │     │  API (Rust)  │     │  (Vectors)  │
└─────────────┘     └──────┬───────┘     └─────────────┘
                           │         ┌─────────────┐
                    ┌──────▼──────┐  │  WebSocket  │
                    │ PostgreSQL  │  │  (Real-Time) │
                    │ (Metadata)  │  └─────────────┘
                    └─────────────┘""", font_size=13)

add_card(slide, 8.3, 2.0, 4.2, 1.2, "🦀 Rust/Axum Backend",
    "30MB Docker image, sub-ms routing\nZero GC pauses, memory safe")
add_card(slide, 8.3, 3.4, 4.2, 1.2, "🔒 Tenant Isolation",
    "Separate Qdrant collections per tenant\nRow-level Postgres isolation\nHashed API keys, encrypted at rest")

add_bullet_list(slide, 0.8, 5.2, 11, 1.5, [
    "SDKs: Python (mrmemory) + TypeScript (memorymr) — published on PyPI + npm",
    "Flexible LLM: BYOK (bring your own key) or managed GPT-4o-mini",
    "Self-hostable: Docker Compose with Postgres + Qdrant + Rust API",
    "0 compile errors, 0 warnings, 38 API tests + 25 SDK tests passing",
], font_size=14)

# ============================================================
# SLIDE 5: Live Demo Reel (Money Slide)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "Live Demo Reel", font_size=40, color=WHITE, bold=True)
add_text(slide, 0.8, 1.8, 11, 0.5, "Features that would take quarters to build — already shipping today.", font_size=18, color=ACCENT2)

add_card(slide, 0.8, 2.5, 3.7, 2.8, "✨ Auto-Remember",
    'Raw conversation in →\nStructured memories out\n\nauto_remember([\n  {"role": "user",\n   "content": "I love Rust"},\n])\n\n→ 2 extracted, 0 duplicates')
add_card(slide, 4.8, 2.5, 3.7, 2.8, "🗜️ Compress",
    "Before: 50 memories\nAfter: 28 memories\n\n44% reduction\nAll meaning preserved\n\nTokens saved per recall:\n~2,200 → ~1,200")
add_card(slide, 8.8, 2.5, 3.7, 2.8, "⚡ Real-Time Sync",
    "Agent A: update(mem_id,\n  content='new info')\n\nAgent B receives:\n  event: memory.updated\n  memory_id: mem_abc\n  content: 'new info'\n\nLatency: <50ms")

add_code_block(slide, 0.8, 5.7, 11.7, 1.2,
    """# LangGraph integration in <10 lines
from mrmemory.langchain import MrMemoryCheckpointer, MrMemoryStore
checkpointer = MrMemoryCheckpointer(api_key="amr_sk_...")
store = MrMemoryStore(api_key="amr_sk_...")
graph = StateGraph(AgentState).compile(checkpointer=checkpointer, store=store)""", font_size=13)

# ============================================================
# SLIDE 6: Differentiation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "Differentiation", font_size=40, color=WHITE, bold=True)
add_text(slide, 0.8, 1.8, 10, 0.5, "Not just a vector store. Not just an open-source framework.", font_size=18, color=ACCENT2)

# Comparison table
headers = ["Feature", "MrMemory", "Mem0", "Letta (ex-MemGPT)", "Raw Qdrant"]
rows = [
    ["Managed API", "✅", "✅", "❌ Self-host", "✅"],
    ["LLM Auto-Extract", "✅", "✅", "✅", "❌"],
    ["Compression", "✅", "❌", "❌", "❌"],
    ["Self-Edit (update/merge/prune)", "✅", "❌", "Partial", "❌"],
    ["Real-Time WebSocket", "✅", "❌", "❌", "❌"],
    ["LangGraph Native", "✅", "❌", "❌", "❌"],
    ["Multi-Agent Sharing", "✅", "❌", "❌", "❌"],
    ["Self-Hostable", "✅", "❌", "✅", "✅"],
    ["Rust Backend", "✅", "❌", "❌ (Python)", "✅"],
]

y_start = 2.5
row_h = 0.38
col_widths = [3.2, 1.8, 1.4, 2.0, 1.8]
x_start = 0.8

# Header row
x = x_start
for i, h in enumerate(headers):
    add_text(slide, x, y_start, col_widths[i], row_h, h, font_size=13, color=ACCENT2, bold=True)
    x += col_widths[i]

# Data rows
for ri, row in enumerate(rows):
    y = y_start + 0.45 + ri * row_h
    x = x_start
    # Subtle alternating bg
    if ri % 2 == 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_start), Inches(y), Inches(sum(col_widths)), Inches(row_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG2
        shape.line.fill.background()
    for ci, cell in enumerate(row):
        c = GREEN if cell == "✅" else (RGBColor(0xEF, 0x44, 0x44) if cell == "❌" else GRAY)
        add_text(slide, x, y, col_widths[ci], row_h, cell, font_size=12, color=c)
        x += col_widths[ci]

# ============================================================
# SLIDE 7: Market Timing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "Market Timing", font_size=40, color=WHITE, bold=True)
add_text(slide, 0.8, 1.8, 10, 0.5, "Memory layer is one of the hottest acquisition categories in AI infra.", font_size=18, color=ACCENT2)

add_card(slide, 0.8, 2.6, 3.7, 2.2, "🦜 LangChain",
    "Raised at $1.25B valuation\nLangGraph = state management\nMissing: intelligent memory layer\nLangSmith needs persistence")
add_card(slide, 4.8, 2.6, 3.7, 2.2, "🔮 Qdrant",
    "Recent $50M Series B\nVector search infrastructure\nMissing: agent intelligence layer\nNeeds managed memory product")
add_card(slide, 8.8, 2.6, 3.7, 2.2, "🌲 Pinecone",
    "Valued at $750M+\nServerless vector DB\nMissing: memory orchestration\nNeeds agent-native features")

add_bullet_list(slide, 0.8, 5.2, 11, 1.5, [
    "Agent engineering is becoming the new platform layer — every AI company needs memory",
    "Vector search infra is exploding but raw vectors ≠ intelligent memory",
    "Mem0 validated the category with $24M total funding — but lacks key features",
    "Window: 6-12 months before memory becomes a commodity feature in every framework",
], font_size=15)

# ============================================================
# SLIDE 8: Founder Velocity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "Founder Velocity & Execution", font_size=40, color=WHITE, bold=True)
add_text(slide, 0.8, 1.8, 10, 0.5, "Solo founder shipping at the speed of larger teams.", font_size=18, color=ACCENT2)

timeline = [
    "Mar 12 — Project conceived, architecture designed",
    "Mar 15 — Rust API scaffold, auth, CRUD, migrations",
    "Mar 17 — Python + TypeScript SDKs built and published",
    "Mar 18 — Backend deployed to Fly.io, semantic search working",
    "Mar 20 — Stripe billing live, auto-provisioning",
    "Mar 22 — LLM auto_remember + LangChain checkpointer shipped",
    "Mar 24 — Compression engine + self-edit tools (update/merge/prune)",
    "Mar 26 — Stats endpoint, Docker Compose self-hosting, full docs",
]
add_bullet_list(slide, 0.8, 2.5, 7, 4, timeline, font_size=14, color=GRAY, spacing=Pt(6))

add_card(slide, 8.5, 2.5, 4, 3.5, "📊 Ship Count (14 days)",
    "• Rust API: 16+ source files\n• 12 API endpoints\n• 2 SDKs published (PyPI + npm)\n• 38 API tests passing\n• 25 SDK tests passing\n• Landing page + docs\n• Stripe billing integration\n• Docker self-hosting\n• 0 compile errors, 0 warnings")

add_text(slide, 0.8, 6.3, 11, 0.5, "0 users today = deliberate pre-launch focus on quality. All infrastructure is production-grade.", font_size=14, color=DARK_GRAY)

# ============================================================
# SLIDE 9: Strategic Fit (LangChain default)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "Strategic Fit", font_size=40, color=WHITE, bold=True)

# LangChain version
add_card(slide, 0.8, 2.0, 11.5, 1.8, "🦜 For LangChain / LangSmith",
    "Integrates natively with LangGraph and LangSmith. Adds managed intelligence, compression,\n"
    "and real-time multi-agent sync to your agent platform. Accelerates production readiness\n"
    "with minimal engineering effort.\nMrMemoryCheckpointer is already a drop-in replacement for LangGraph.")

# Qdrant version
add_card(slide, 0.8, 4.1, 11.5, 1.5, "🔮 For Qdrant",
    "Built on Qdrant with Rust backend. Extends your vector engine into a full intelligent memory\n"
    "layer with auto-orchestration, compression, self-edits, and live collaboration.\n"
    "Rust-native extension of your vector engine — zero rewrite needed.")

# Pinecone version
add_card(slide, 0.8, 5.9, 11.5, 1.2, "🌲 For Pinecone",
    "Transforms long-term vector memory into complete intelligent agent memory. Adds auto_remember,\n"
    "compression, self-edit tools, and real-time sync.\nInstant upgrade from long-term vectors to full agent memory.")

add_text(slide, 0.8, 7.0, 11, 0.3, "Customize this slide per buyer. Each version highlights unique strategic alignment.", font_size=12, color=DARK_GRAY)

# ============================================================
# SLIDE 10: The Ask
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "The Ask", font_size=40, color=WHITE, bold=True)

add_card(slide, 0.8, 2.2, 5.5, 3.5, "💰 Acqui-Hire Terms",
    "Acqui-hire range: $1.2M – $2.5M\n(cash + equity refresh + founder role\nat market rate)\n\n"
    "Timeline: 4–8 week integration\n\n"
    "Role: Founding Engineer or\nMemory Product Lead\n\n"
    "Value: Immediate production-grade\nmemory layer + founder who ships fast")

add_card(slide, 6.8, 2.2, 5.5, 3.5, "🎁 What You Get",
    "✅ Production Rust API (0 errors, 0 warnings)\n"
    "✅ 2 published SDKs (Python + TypeScript)\n"
    "✅ LangGraph-native integration\n"
    "✅ Auto-remember + Compression + Self-Edit\n"
    "✅ Real-time WebSocket multi-agent sync\n"
    "✅ Stripe billing, Docker self-hosting\n"
    "✅ Full docs, landing page, test suites\n"
    "✅ A founder who ships fast — solo")

add_text(slide, 0.8, 6.0, 11, 0.5, "This is months of runway compressed into 14 days of execution.", font_size=16, color=ACCENT2)

# ============================================================
# SLIDE 11: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_accent_line(slide, 0.8, 0.8, 1.5)
add_text(slide, 0.8, 1.0, 11, 0.8, "Next Steps", font_size=40, color=WHITE, bold=True)

add_bullet_list(slide, 0.8, 2.2, 10, 3, [
    "📞  Live 15-minute code + architecture walkthrough under NDA",
    "🔓  Immediate access to full repo, private backend, and SDKs",
    "🤝  Let's make MrMemory the default memory backend",
    "",
    "Ready for NDA, full repo access, and a 15-minute technical walkthrough.",
], font_size=20, color=GRAY, spacing=Pt(16))

add_text(slide, 0.8, 5.0, 11, 0.8, "Darren R", font_size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 5.7, 11, 0.4, "Solo Founder, MrMemory", font_size=18, color=ACCENT2)

add_text(slide, 0.8, 6.3, 5, 0.8,
    "🌐  mrmemory.dev\n📦  github.com/masterdarren23/mrmemory",
    font_size=14, color=GRAY)
add_text(slide, 6, 6.3, 5, 0.8,
    "🐍  pypi.org/project/mrmemory\n📦  npmjs.com/package/memorymr",
    font_size=14, color=GRAY)

add_text(slide, 3, 7.0, 7, 0.4, "Thank You 🧠", font_size=24, color=ACCENT2, alignment=PP_ALIGN.CENTER, bold=True)

# Save
out = r"C:\Users\johnl\.openclaw\workspace\amr-project\pitch-deck\MrMemory-Pitch-Deck.pptx"
prs.save(out)
print(f"Saved to {out}")
